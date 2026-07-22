import os
import re
import json
import io
import pandas as pd
from docx import Document
from pptx import Presentation

# Avoid pypdf import crash if not fully installed yet (load dynamically)
def extract_pdf_text(file_bytes):
    try:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(file_bytes))
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text
    except Exception as e:
        return f"Error extracting PDF: {str(e)}"

# Domain detection function based on keywords
def detect_domain(text, df=None):
    text = text.lower()
    
    # Also look at DataFrame columns if available
    col_str = ""
    if df is not None:
        col_str = " ".join([str(c) for c in df.columns]).lower()
    
    combined = text + " " + col_str
    
    domains = {
        "Retail": [
            "sales", "quantity", "price", "store", "customer", "transaction", "revenue", 
            "item", "product", "sku", "discount", "order", "inventory", "salesperson", "retail"
        ],
        "Finance": [
            "assets", "liabilities", "equity", "revenue", "profit", "ebitda", "quarter", 
            "portfolio", "rate", "yield", "balance", "expense", "stock", "dividend", "interest", "finance"
        ],
        "Healthcare": [
            "patient", "doctor", "diagnosis", "admission", "treatment", "prescription", 
            "hospital", "billing", "clinical", "symptom", "disease", "healthcare", "ehr", "medical"
        ],
        "Education": [
            "student", "grade", "gpa", "score", "course", "term", "attendance", 
            "enrollment", "teacher", "school", "assignment", "test", "tuition", "education"
        ],
        "Manufacturing": [
            "production", "yield", "defect", "downtime", "machine", "sensor", "part", 
            "assembly", "batch", "maintenance", "oee", "manufacturing", "factory", "line"
        ]
    }
    
    scores = {d: 0 for d in domains}
    for dom, keywords in domains.items():
        for kw in keywords:
            # Count word matches
            matches = len(re.findall(r'\b' + re.escape(kw) + r'\b', combined))
            # Substring matches as fallback
            if matches == 0 and kw in combined:
                matches = 1
            scores[dom] += matches
            
    best_domain = max(scores, key=scores.get)
    if scores[best_domain] == 0:
        return "General Business"
    return best_domain

# SQL Dump Parser
def parse_sql_dump(sql_text):
    # Try to find INSERT INTO statements
    # e.g., INSERT INTO `table` (`col1`, `col2`) VALUES (1, 'val'), (2, 'val2');
    table_data = {}
    
    insert_pattern = re.compile(
        r"INSERT\s+INTO\s+[`\"']?(\w+)[`\"']?\s*(?:\(([^)]+)\))?\s*VALUES\s*(.+?)(?:;|$)", 
        re.IGNORECASE | re.DOTALL
    )
    
    for match in insert_pattern.finditer(sql_text):
        table_name = match.group(1)
        columns_raw = match.group(2)
        values_raw = match.group(3)
        
        # Clean columns
        cols = []
        if columns_raw:
            cols = [c.strip(" `\"'") for c in columns_raw.split(",")]
            
        # Parse values
        # e.g., (1, 'val', 3.5), (2, 'val2', NULL)
        row_pattern = re.compile(r"\((.+?)\)(?:\s*,\s*|\s*$)")
        rows = []
        for row_match in row_pattern.finditer(values_raw):
            row_content = row_match.group(1)
            # Parse commas not inside strings
            # Simple regex parser for values: matches strings '...', "..." or numbers/NULL
            val_items = []
            for item in re.split(r",(?=(?:[^']*'[^']*')*[^']*$)", row_content):
                item = item.strip()
                if item.upper() == "NULL":
                    val_items.append(None)
                elif (item.startswith("'") and item.endswith("'")) or (item.startswith('"') and item.endswith('"')):
                    val_items.append(item[1:-1])
                else:
                    # Try to convert to float or int
                    try:
                        if "." in item:
                            val_items.append(float(item))
                        else:
                            val_items.append(int(item))
                    except ValueError:
                        val_items.append(item)
            rows.append(val_items)
            
        if not cols and rows:
            # Generate default column names if not specified in SQL
            cols = [f"column_{i+1}" for i in range(len(rows[0]))]
            
        if table_name not in table_data:
            table_data[table_name] = {"cols": cols, "rows": []}
        
        # Make sure column lengths match rows
        table_data[table_name]["rows"].extend(rows)
        
    return table_data

import numpy as np
import datetime
import math

def sanitize_value(v):
    if v is None:
        return None
    
    # Pandas / Numpy NA / NaT checks
    try:
        if pd.isna(v):
            return None
    except Exception:
        pass
        
    if isinstance(v, (float, np.floating)):
        if math.isnan(v) or math.isinf(v):
            return None
        return float(v)
    elif isinstance(v, (int, np.integer)):
        return int(v)
    elif isinstance(v, (bool, np.bool_)):
        return bool(v)
    elif isinstance(v, (datetime.datetime, datetime.date, pd.Timestamp)):
        return v.isoformat()
    elif isinstance(v, (list, tuple)):
        return [sanitize_value(x) for x in v]
    elif isinstance(v, dict):
        return {str(rk): sanitize_value(rv) for rk, rv in v.items()}
    else:
        s = str(v).strip()
        if s.lower() in ("nan", "none", "null", "nat", "<na>"):
            return None
        return s

def clean_dataframe_to_dict(df):
    if df is None or df.empty:
        return []
    
    df = df.copy()
    
    # 1. Clean and deduplicate column names
    cleaned_cols = []
    seen = {}
    for i, col in enumerate(df.columns):
        col_str = str(col).strip() if col is not None else ""
        if not col_str or col_str.lower().startswith("unnamed:"):
            col_str = f"Column_{i+1}"
        if col_str in seen:
            seen[col_str] += 1
            col_str = f"{col_str}_{seen[col_str]}"
        else:
            seen[col_str] = 1
        cleaned_cols.append(col_str)
    
    df.columns = cleaned_cols
    
    # 2. Extract records safely with JSON-compatible types
    records = []
    raw_dict = df.to_dict(orient="records")
    for row in raw_dict:
        cleaned_row = {}
        for k, v in row.items():
            cleaned_row[str(k)] = sanitize_value(v)
        records.append(cleaned_row)
        
    return records

# Master Parser function
def parse_file(filename, file_bytes):
    ext = os.path.splitext(filename)[1].lower()
    df = None
    domain = "General Business"
    doc_type = "Structured Dataset"
    text_content = ""
    
    # 1. Excel files
    if ext in [".xlsx", ".xls", ".xlsm", ".xlsb", ".xltx", ".ods"]:
        doc_type = "Excel Spreadsheet"
        try:
            # Try engines sequentially
            for engine in [None, "openpyxl", "xlrd", "pyxlsb", "odf"]:
                try:
                    if engine:
                        df = pd.read_excel(io.BytesIO(file_bytes), engine=engine, nrows=10000)
                    else:
                        df = pd.read_excel(io.BytesIO(file_bytes), nrows=10000)
                    if df is not None and not df.empty:
                        break
                except Exception:
                    continue
                    
            # Direct openpyxl fallback if pandas read_excel failed
            if df is None or df.empty:
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
                    sheet = wb.active
                    data = list(sheet.values)
                    if data and len(data) > 1:
                        header = [str(c) if c is not None else f"Column_{i+1}" for i, c in enumerate(data[0])]
                        rows = data[1:]
                        df = pd.DataFrame(rows, columns=header)
                except Exception:
                    pass
                    
            if df is None or df.empty:
                raise Exception("No non-empty tables could be read from this Excel file.")
                
            text_content = df.head(10).to_string()
            domain = detect_domain(text_content, df)
        except Exception as e:
            raise Exception(f"Failed to parse Excel file: {str(e)}")
            
    # 2. CSV files
    elif ext in [".csv", ".tsv", ".txt", ".dat", ".log"]:
        doc_type = "Comma Separated Values"
        try:
            for enc in ["utf-8", "latin-1", "cp1252", "utf-16"]:
                try:
                    for sep in [None, ",", "\t", "|", ";"]:
                        try:
                            if sep:
                                df = pd.read_csv(io.BytesIO(file_bytes), encoding=enc, sep=sep, engine="python", nrows=10000)
                            else:
                                df = pd.read_csv(io.BytesIO(file_bytes), encoding=enc, nrows=10000)
                            if df is not None and not df.empty and len(df.columns) > 0:
                                break
                        except Exception:
                            continue
                    if df is not None and not df.empty:
                        break
                except Exception:
                    continue
                    
            if df is None or df.empty:
                lines = file_bytes.decode("utf-8", errors="ignore").splitlines()
                lines = [l.strip() for l in lines if l.strip()]
                df = pd.DataFrame({"line": range(1, len(lines)+1), "content": lines})
                
            text_content = df.head(10).to_string()
            domain = detect_domain(text_content, df)
        except Exception as e:
            raise Exception(f"Failed to parse CSV/Text file: {str(e)}")
            
    # 3. JSON files
    elif ext == ".json":
        doc_type = "JSON Document"
        try:
            data = json.loads(file_bytes.decode('utf-8'))
            if isinstance(data, list):
                df = pd.DataFrame(data)
            elif isinstance(data, dict):
                # Try to normalize or read keys
                if any(isinstance(v, list) for v in data.values()):
                    # Find list values
                    list_key = [k for k, v in data.items() if isinstance(v, list)][0]
                    df = pd.DataFrame(data[list_key])
                else:
                    df = pd.DataFrame([data])
            else:
                df = pd.DataFrame([{"value": data}])
            text_content = json.dumps(data)[:2000]
            domain = detect_domain(text_content, df)
        except Exception as e:
            raise Exception(f"Failed to parse JSON file: {str(e)}")
            
    # 4. Word Documents (.docx)
    elif ext == ".docx":
        doc_type = "Word Document"
        try:
            doc = Document(io.BytesIO(file_bytes))
            text_content = "\n".join([p.text for p in doc.paragraphs])
            domain = detect_domain(text_content)
            
            # Extract tables if any exist in the word doc
            tables_data = []
            for t_idx, table in enumerate(doc.tables):
                t_rows = []
                for row in table.rows:
                    t_rows.append([cell.text.strip() for cell in row.cells])
                if t_rows:
                    cols = t_rows[0]
                    rows = t_rows[1:]
                    # Filter empty headers
                    cols = [c if c else f"col_{i}" for i, c in enumerate(cols)]
                    tdf = pd.DataFrame(rows, columns=cols)
                    tables_data.append(tdf)
            
            if tables_data:
                df = tables_data[0] # Use the first table as main dataframe
            else:
                # Create a simple unstructured dataset from paragraphs
                lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
                df = pd.DataFrame({"paragraph_index": range(1, len(lines)+1), "content": lines})
        except Exception as e:
            raise Exception(f"Failed to parse Word Document: {str(e)}")
            
    # 5. PowerPoint Presentation (.pptx)
    elif ext == ".pptx":
        doc_type = "PowerPoint Presentation"
        try:
            prs = Presentation(io.BytesIO(file_bytes))
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_runs.append(shape.text.strip())
            text_content = "\n".join(text_runs)
            domain = detect_domain(text_content)
            
            lines = [line for line in text_runs if line]
            df = pd.DataFrame({"slide_index": range(1, len(lines)+1), "content": lines})
        except Exception as e:
            raise Exception(f"Failed to parse PowerPoint presentation: {str(e)}")
            
    # 6. PDF files
    elif ext == ".pdf":
        doc_type = "PDF Document"
        text_content = extract_pdf_text(file_bytes)
        domain = detect_domain(text_content)
        
        # If there's structured text, try to extract tabular data or split text
        lines = [line.strip() for line in text_content.split("\n") if line.strip()]
        if len(lines) > 5 and ("," in lines[0] or "\t" in lines[0] or "  " in lines[0]):
            # Simple text parsing try
            try:
                from io import StringIO
                df = pd.read_csv(StringIO(text_content), sep=None, engine='python')
            except Exception:
                pass
        
        if df is None or df.empty or len(df.columns) <= 1:
            # Fallback to lines of text
            df = pd.DataFrame({"page_line": range(1, len(lines)+1), "content": lines})
            
    # 7. SQL Dump
    elif ext in [".sql"]:
        doc_type = "SQL Database Dump"
        try:
            sql_text = file_bytes.decode('utf-8', errors='ignore')
            text_content = sql_text[:2000]
            tables = parse_sql_dump(sql_text)
            
            if tables:
                # Pick the table with the most rows as primary
                primary_table = max(tables.keys(), key=lambda k: len(tables[k]["rows"]))
                cols = tables[primary_table]["cols"]
                rows = tables[primary_table]["rows"]
                df = pd.DataFrame(rows, columns=cols)
                domain = detect_domain(text_content, df)
            else:
                lines = [line.strip() for line in sql_text.split("\n") if line.strip()]
                df = pd.DataFrame({"line": range(1, len(lines)+1), "statement": lines})
                domain = "General Business"
        except Exception as e:
            raise Exception(f"Failed to parse SQL file: {str(e)}")
            
    # 8. Images (or Scanned Documents)
    elif ext in [".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".webp", ".gif", ".jfif", ".heic", ".heif", ".svg", ".ico", ".avif"]:
        doc_type = "Scanned Image (OCR)"
        try:
            from backend.ocr import perform_ocr
            df, domain, text_content = perform_ocr(filename, file_bytes)
        except Exception:
            df = None
            domain = "Retail"
            text_content = "Image Asset Processed"
        
    else:
        # Universal fallback for any other text/binary document
        doc_type = f"Document ({ext.upper() if ext else 'DATA'})"
        try:
            text_content = file_bytes.decode('utf-8', errors='ignore')
            lines = [l.strip() for l in text_content.split("\n") if l.strip()]
            if lines:
                first_line = lines[0]
                sep = "," if "," in first_line else ("\t" if "\t" in first_line else ("|" if "|" in first_line else (";" if ";" in first_line else None)))
                if sep and len(lines) > 1:
                    try:
                        df = pd.read_csv(io.StringIO(text_content), sep=sep, nrows=10000)
                    except Exception:
                        df = pd.DataFrame({"line": range(1, len(lines)+1), "content": lines})
                else:
                    df = pd.DataFrame({"line": range(1, len(lines)+1), "content": lines})
            else:
                df = pd.DataFrame({"info": ["File parsed successfully but contains no readable content."]})
            domain = detect_domain(text_content[:2000], df)
        except Exception:
            df = pd.DataFrame({"info": [f"Raw file asset loaded ({len(file_bytes)} bytes)."]})
            domain = "General Business"
        
    return df, domain, doc_type, text_content
